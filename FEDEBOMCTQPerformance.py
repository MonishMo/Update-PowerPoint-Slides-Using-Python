# -*- coding: utf-8 -*-
"""
@author- Andrew Haakenson (AHAAKENS)  & Monish (BMONISH)

Automatically updates the FEDEBOM CTQ Performance slide with global averages for each CTQ  of the Business Deck.
"""

import pptx
import pandas as pd
import numpy as np
import calendar

import ctq_query
import model


def find_threshold(ctq,red):
    '''Find and remove the Outlier then Calculates the Global Average based on the thersold value for each CTQs
    
    Parameters: {ctq: CTQ NAME,
                 red : Global Red Value (Upper Thresold)}
    
    Return: Average Glabal Average Without Outliers(Float)'''
    
    #Fetches Query String from ctq_query file for selected CTQ
    query = ctq_query.fetch_query(ctq)

    param = ctq_query.get_parameter()
    
    oneshot_search = model.oneshot_search(query,param)
    oneshot_search = pd.DataFrame(oneshot_search)
    
    oneshot_search['Duration [sec]'] = pd.to_numeric(oneshot_search['Duration [sec]'])
    
    result = {}
    count = 0
    total_records = len(oneshot_search)
    #user_name = []
    user_id_country = []
    for index in range(len(oneshot_search)):
        if ((oneshot_search['Duration [sec]'][index:].mean() < red)):
            #Removes biggest value one by one and checks the average
            #Once average is less than Red, fetch the average and return 
            list_without_outliers = oneshot_search['Duration [sec]'][index:].to_list()
            
            break
        #user_name.append(oneshot_search['CDSID'][index])
        count+=1
        user_id_country.append(oneshot_search['CDSID'][index]+"["+oneshot_search['Country'][index]+"]")
    average_without_outliers = sum(list_without_outliers)/len(list_without_outliers)

    average_without_outliers = np.float64(round(average_without_outliers,2))
    
    result['total_records'] = total_records
    result['count'] = count
    result['average'] = average_without_outliers
    #result['user_name'] = user_name
    result['user_id_country'] = ",".join(user_id_country)
    
    return result

def update_slide(presentation_file_name): 
    '''Updates Global Average column on FEDEBOM CTQ Performance slide. Data is fetched from Splunk through Splunk Enterprise SDK for Python
    Values are classified with Colors (Red,Yellow,Green)
    
    Return: Modifies the existing powerpoint slide with latest data'''
    
    query = ctq_query.fetch_query("CTQ Daily Status Dashboard")

    param = ctq_query.get_parameter()

    ctq_dashboard = model.oneshot_search(query,param)
    ctq_dashboard = pd.DataFrame(ctq_dashboard)
    ctq_dashboard.set_index('CTQ',inplace=True)

    # Color constants
    green = pptx.dml.color.RGBColor(74, 161, 53)
    yellow = pptx.dml.color.RGBColor(255, 255, 0)
    #red = pptx.dml.color.RGBColor(255, 0, 0)

    '''
    Presentation file name
    NOTE: Change this to current presentation file name
    '''
    #presentation_file_name = "Eng BOM Prod Sup And CTQ Review.pptx"

    prs = pptx.Presentation(presentation_file_name)

    # Select the table on slide 2
    table = prs.slides[1].shapes[3].table

    table_row = 1

    # Extract each row from the search result and insert value into the table
    #for line in reader:
    for index in range(1,len(ctq_dashboard)):
        # Stop if it goes past the end of the table
        if table_row > 46:
            break

        ctq = table.cell(table_row,0).text
        cell = table.cell(table_row,2)
        cell.text = ctq_dashboard.loc[ctq]['Global']

        # Extract the threshold values from the table
        # For the green value, skip the '<=' sign at the beginning
        lower_threshold = float(table.cell(table_row, 1).text[2:]) 
        middle_threshold = float(table.cell(table_row + 1, 1).text)
        upper_threshold = float(table.cell(table_row + 2, 1).text)


        ctq_average = float(ctq_dashboard.loc[ctq]['Global'])
        # Determine the coloring of each cell
        if ctq_average > upper_threshold:
            #find the average global count after removing outliers
            result = find_threshold(ctq,upper_threshold)
            
            without_outlier_cell = table.cell(table_row,2)
            without_outlier_cell.text = str(result['average'])
            
            comment = """%d of %d transactions had more than 106000 secs of duration time which are usually due to client application that are not shutdown properly. Users affected: %s""" %(result['count'],result['total_records'],result['user_id_country'])
            
            cell = table.cell(table_row,3)
            cell.text = ctq_dashboard.loc[ctq]['Global']
            
            comment_cell = table.cell(table_row,4)
            comment_cell.text = comment
            
            if result['average'] > lower_threshold:
                color = yellow
            else:
                color = green
                
                        
            without_outlier_cell.fill.solid()
            without_outlier_cell.fill.fore_color.rgb = color
            
            paragraph = without_outlier_cell.text_frame.paragraphs[0]
            paragraph.font.size = pptx.util.Pt(9)
            paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = pptx.util.Pt(9)
            paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER
            
            paragraph = comment_cell.text_frame.paragraphs[0]
            paragraph.font.size = pptx.util.Pt(9)
            paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.LEFT
            
            table_row += 3
            
            continue 
        
        elif ctq_average > middle_threshold:
            color = yellow
        else:
            color = green
            
        outlier_cell = table.cell(table_row,3)
        outlier_cell.text = ""
        comment_cell = table.cell(table_row,4)
        comment_cell.text = ""

        cell.fill.solid()
        cell.fill.fore_color.rgb = color

        # Fix the styling to be 9 pt font and centered
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = pptx.util.Pt(9)
        paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER

        # Skip 3 rows down due to merged rows
        table_row += 3
        
    param_dict = ctq_query.get_parameter()
    latest_time = param_dict["latest_time"]
    earliest_time = param_dict["earliest_time"]
    
    latest_date = latest_time.split("T")[0]
    earliest_date = earliest_time.split("T")[0]
    
    '''month = int(latest_time.split("-")[1])
    date = int(latest_time.split("-")[2][:2])
    year = int(latest_time.split("-")[0])
    
    month_name = calendar.month_name[month]'''
    
    slide = prs.slides[1]
    
    title = slide.shapes[1]
    title.text = "FEDEBOM Performance (KPIs/CTQs) ({0} to {1})".format(earliest_date,latest_date)
    title.text_frame.fit_text()
    #note_slide = slide.notes_slide
    #notes = note_slide.notes_text_frame
    
    #notes.text = "Note : Metrics are until %d %s %d"%(date,month_name,year)


    prs.save(presentation_file_name)
    print("FEDEBOM Performance (KPIs/CTQs) Slide Updated Successfully")
