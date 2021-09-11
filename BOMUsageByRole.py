# -*- coding: utf-8 -*-
"""
@author- Andrew Haakenson (ahaakens) & Monish

Automatically updates both charts on Monthly BOM Usage by Role slide of the Business Deck.

Data is extracted from Splunk using Splunk Enterprise SDK.

Please change the presentation file name to the current one in the constants
section before running.
"""

import ctq_query
import model

from pptx import Presentation
from pptx.chart.data import CategoryChartData
import pandas as pd
import calendar

'''#Converting result(List) into DataFrame to easily work.
bom_search = pd.DataFrame(bom_search)
bom_save_submit = pd.DataFrame(bom_save_submit)'''

def update_chart(df, column2_header, chart_num,presentation_file_name, slide4_shape_offset):
    '''Gets the result from the Splunk and update the Chart Using pptx package
    Parameters: {df : DataFrame contains Result from Splunk,
                 column2_header: Specifies Header (Search or Submit)
                 chart_num: Placeholder on slide to update the chart'''
    
    # Dictionary to keep track of users by role
    users = {}
    # Dictionary to keep track of number of searches for each role
    searches_by_role = {}
    
    # Parse the file and extract the necessary columns
        
    for index,row in df.iterrows():
        user = row['User']
        role = row['Role']  
        
        # Skip this entry if no user role recorded
        if role == "No User Role Captured":
            continue
        
        # Role already encountered and in dictionaries
        if role in searches_by_role.keys():
            searches_by_role[role] += 1
            if user not in users[role]:
                users[role][user] = True
        
        else:
            searches_by_role[role] = 1
            users[role] = {}
            users[role][user] = True
            
    # 3 lists, one for each column
    # first list is for user roles, second is for total users, third for total searches
    split_file = [[], [], []] 
    
    split_file[0] = list(searches_by_role.keys())
    # Sort roles alphabetically
    split_file[0].sort()
    for role in split_file[0]:
        split_file[2].append(searches_by_role[role])
        split_file[1].append(len(users[role].keys()))
    
    # Load the powerpoint from a pptx file
    prs = Presentation(presentation_file_name)
    
    chart_data = CategoryChartData()
    chart_data.categories = split_file[0] # Add the user roles on the x axis
    chart_data.add_series("Total Users", split_file[1]) # Add the total users
    chart_data.add_series(column2_header, split_file[2]) # Add the total searches
    
    chart = prs.slides[3].shapes[chart_num + slide4_shape_offset].chart
    # Replace the current data with data from csv file
    # Only replaces data. Chart formatting remains unchanged
    chart.replace_data(chart_data)
    
    param_dict = ctq_query.get_parameter()
    latest_time = param_dict["latest_time"]
    earliest_time = param_dict["earliest_time"]
    
    latest_date = latest_time.split("T")[0]
    earliest_date = earliest_time.split("T")[0]
    
    '''month = int(latest_time.split("-")[1])
    date = int(latest_time.split("-")[2][:2])
    year = int(latest_time.split("-")[0])
    
    month_name = calendar.month_name[month]'''
    
    slide = prs.slides[3]
    
    title = slide.shapes[1]
    title.text = "Monthly BOM Usage by Role ({0} to {1})".format(earliest_date,latest_date)
    title.text_frame.fit_text()
    
    #note_slide = slide.notes_slide
    #notes = note_slide.notes_text_frame
    
    #notes.text = "Note : Metrics are until %d %s %d"%(date,month_name,year)
    
    prs.save(presentation_file_name)
    
def bom_usage(presentation_file_name):
    """Function that calls the Main method to update the chart in Powerpoint Slide"""
    
    # Constants
    # How many shapes are in front of charts on slide 4 (add this to chart number
    # to get shape number)
    slide4_shape_offset = 1
    
    '''
    File names
    NOTE: Change presentation file name to current file name
    '''
    #presentation_file_name = "Eng BOM Prod Sup And CTQ Review.pptx"
    
    #Fetch the Query String BOMSearchAndLoadworksapce and BOMSaveSubmitChangesProduct     
    bom_search_query = ctq_query.fetch_query("BOMSearchAndLoadworksapce")
    bom_save_submit_query = ctq_query.fetch_query("BOMSaveSubmitChangesProduct")
    
    param = ctq_query.get_parameter()
    
    #Result from Splunk Search Query
    bom_search = model.oneshot_search(bom_search_query,param)
    bom_save_submit = model.oneshot_search(bom_save_submit_query,param)
    bom_search = pd.DataFrame(bom_search)
    bom_save_submit = pd.DataFrame(bom_save_submit)
    
    update_chart(bom_search, "Total Searches", 1, presentation_file_name, slide4_shape_offset)
    update_chart(bom_save_submit, "Total Saves", 2, presentation_file_name, slide4_shape_offset)
    
    print("BOM Usage By Role Slide updated Successfully")
