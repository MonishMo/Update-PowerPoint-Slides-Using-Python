# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""


import cx_Oracle
from getpass import getpass

from datetime import datetime
import pptx
import pandas as pd
import calendar

import model
import ctq_query

def get_average_count(metrics):
    
    '''calculate the average count for each metrics based the Splunk input search query
    
    Parameter: Name of the metrics
    
    Return: Average Count for each metrics(returned with str type for updating the value in pptx)
    '''
    
    param = ctq_query.get_parameter()
    
    query = ctq_query.fetch_query(metrics)
    usability_count = model.oneshot_search(query,param)
    usability_count = pd.DataFrame(usability_count)
    
    count = [[],[]]
    
    #skip US holidays 
    with open("Holidays.txt","r") as holiday_file:
        line = holiday_file.read()
        holidays = line.split()
        
    for index, item in usability_count.iterrows():
        
        date_str = str(item['date'])+"-2021"
        date = datetime.strptime(date_str,'%d-%m-%Y')
        day_of_week = date.weekday()
    
        if day_of_week != 5 and day_of_week != 6:
            if date_str not in holidays:
                count[0].append(item['date'])
                count[1].append(int(item.iloc[0]))
            
        avg_count  = int(sum(count[1])/len(count[1]))
        
    return str(avg_count)
    
def get_registered_user():
    
    '''db_password = getpass("Enter Password for connecting to DB: ")
    
    cx_Oracle.init_oracle_client(lib_dir= r"C:\oracle\instantclient-basic-windows.x64-19.11.0.0.0dbru\instantclient_19_11")
    print("Oracle Client Version ",cx_Oracle.clientversion())
    
    query = ctq_query.fetch_query("Registered_User")
    cdsid_lst = []

    connection = cx_Oracle.connect('BOMINFEDE_READ4/db_password@bomfdb.ford.com:1521/PDPROD10')
    cursor = connection.cursor()
    
    cursor.execute(query)
    for result in cursor:
        cdsid_lst.append(result)
        
    total_user = len(cdsid_lst)'''
    
    total_user = 1201    
    
    return str(total_user)

def update_slide(user_count,presentation_file_name):
    '''Updates Average Registered User Count, Active User, Search Count, Submit Count
    on FEDEBOM User Metrics Slide. 
    
    Data is fetched from Splunk through Splunk Enterprise SDK for Python
    Weekend Values are removed and average is found.
    
    Return: Modifies the existing powerpoint slide with latest data'''
    
    #edit existing pptx
    #presentation_file_name = "Eng BOM Prod Sup And CTQ Review.pptx"

    prs = pptx.Presentation(presentation_file_name)
    
    # Select the table on slide 2
    table = prs.slides[2].shapes[4].table
    
    row = 1
    
    #update value in the table
    for key, value in user_count.items():
        
        new_cell = table.cell(row, 1)
        
        new_cell.text = table.cell(row,2).text
        
        old_cell = table.cell(row,2)
        old_cell.text = user_count[key]
        
        row+=1
        
        #changing the font size and alignment 
        paragraph = new_cell.text_frame.paragraphs[0]
        paragraph.font.size = pptx.util.Pt(9)
        paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER
        
        paragraph = old_cell.text_frame.paragraphs[0]
        paragraph.font.size = pptx.util.Pt(9)
        paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER
        
        prs.save(presentation_file_name)
        
    param_dict = ctq_query.get_parameter()
    
    latest_date = param_dict["latest_time"].split("T")[0].split("-")
    earliest_date = param_dict["earliest_time"].split("T")[0].split("-")
    
    #changing the date format to MM/DD/YYY
    latest_month = int(latest_date[1])
    latest_date_num = int(latest_date[2])
    latest_year = int(latest_date[0])
    
    earliest_month = int(earliest_date[1])
    earliest_date_num = int(earliest_date[2])
    earliest_year = int(earliest_date[0])
    
    latest_date_format = "{0}/{1}/{2}".format(latest_month,latest_date_num,latest_year)
    earliest_date_format = "{0}/{1}/{2}".format(earliest_month,earliest_date_num,earliest_year)
    #month_name = calendar.month_name[month]
    
    slide = prs.slides[2]
    
    title = slide.shapes[0]
    title.text = "FEDEBOM User Metrics ({0} to {1})".format(earliest_date_format,latest_date_format)
    title.text_frame.fit_text()
    
    #note_slide = slide.notes_slide
    #notes = note_slide.notes_text_frame
    
    #notes.text = "Note : Metrics are until %d %s %d"%(date,month_name,year)
        
    prs.save(presentation_file_name)
    
def user_metrics(presentation_file_name):
    
    #avg_user_count = update_table(con)
    
    user_count ={}
    
    registered_user = get_registered_user()
    user_count['registered_user'] = registered_user
    
    active_user = get_average_count("ActiveUser")
    user_count['active_user'] = active_user
    
    search_count = get_average_count("SearchCount")
    user_count['search_count'] = search_count
    
    submit_count = get_average_count("SubmitCount")
    user_count['submit_count'] = submit_count
    
    print("Average count :\n",user_count)
    
    update_slide(user_count,presentation_file_name)

    print("FEDEBOM User Metrics  updated successfully")